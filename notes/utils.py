# notes/utils.py
from io import BytesIO
from django.http import HttpResponse
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
from docx import Document
from docx.shared import RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor as PptxRGBColor
import re

from bs4 import BeautifulSoup
import google.generativeai as genai

def generate_notes_gemini(topic: str) -> list:
    """
    Generate structured notes and return a list of dicts:
    [{'type': 'heading'/'subheading'/'paragraph'/'note'/'list', 'text': str, 'level': int}]
    """
    model = genai.GenerativeModel("models/gemini-pro-latest")
    prompt = f"""
    Generate detailed, student-friendly notes on the topic: {topic}.
    Return clean HTML with:
    - <h1> for main title
    - <h2> for sections
    - <h3> for sub-sections
    - <p> for normal text
    - <ul><li> for bullet points (nested allowed)
    - <b> for key terms
    - Use <p><b style='color:red'>NOTE:</b> text</p> for notes
    Do NOT use markdown.
    """

    response = model.generate_content(prompt)
    html_content = response.text if response and response.text else "<p>No notes generated.</p>"

    soup = BeautifulSoup(html_content, "html.parser")
    structured_notes = []

    def parse_ul(ul_tag, level=0):
        items = []
        for li in ul_tag.find_all("li", recursive=False):
            text = li.get_text(strip=True)
            items.append({"type": "list", "text": text, "level": level})
            nested_ul = li.find("ul", recursive=False)
            if nested_ul:
                items.extend(parse_ul(nested_ul, level + 1))
        return items

    for tag in soup.find_all(["h1", "h2", "h3", "p", "ul"]):
        if tag.name == "h1":
            structured_notes.append({"type": "heading", "text": tag.get_text(strip=True)})
        elif tag.name == "h2":
            structured_notes.append({"type": "subheading", "text": tag.get_text(strip=True)})
        elif tag.name == "h3":
            structured_notes.append({"type": "subsubheading", "text": tag.get_text(strip=True)})
        elif tag.name == "p":
            if tag.find("b", style=lambda s: s and "color:red" in s):
                structured_notes.append({"type": "note", "text": tag.get_text(strip=True).replace("NOTE:", "").strip()})
            else:
                structured_notes.append({"type": "paragraph", "text": tag.get_text(strip=True)})
        elif tag.name == "ul":
            structured_notes.extend(parse_ul(tag))

    return structured_notes

# ---------------- PDF ----------------
def generate_pdf(topic: str, notes_input) -> HttpResponse:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()

    story = []
    story.append(Paragraph(f"<b>Notes on {topic}</b>", styles["Title"]))
    story.append(Spacer(1, 12))

    # CASE A: structured notes (preferred)
    if isinstance(notes_input, list):
        for item in notes_input:
            typ = item.get('type')
            text = (item.get('text') or "").strip()
            if not text:
                continue

            if typ == "heading":
                story.append(Paragraph(f"<b>{text}</b>", styles["Heading1"]))
            elif typ == "subheading":
                story.append(Paragraph(f"<b>{text}</b>", styles["Heading2"]))
            elif typ == "subsubheading":
                story.append(Paragraph(f"<b>{text}</b>", styles["Heading3"]))
            elif typ == "paragraph":
                story.append(Paragraph(text, styles["Normal"]))
            elif typ == "note":
                story.append(Paragraph(
                    f"<b><font color='red'>NOTE:</font></b> {text}",
                    styles["Normal"]
                ))
            elif typ == "list":
                # indent bullet points based on nesting level
                indent = 12 * (item.get("level", 0) or 0)
                story.append(Paragraph(f"• {text}", styles["Normal"]))
                story[-1].leftIndent = indent
            story.append(Spacer(1, 6))

    # CASE B: fallback to plain text
    else:
        notes_text = (notes_input or "")
        for line in notes_text.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("NOTE:"):
                story.append(Paragraph(f"<b><font color='red'>{line}</font></b>", styles["Normal"]))
            else:
                story.append(Paragraph(line, styles["Normal"]))
            story.append(Spacer(1, 6))

    doc.build(story)
    buffer.seek(0)
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{topic}.pdf"'
    return response


# ---------------- WORD ----------------
def generate_word(topic: str, notes_text: str) -> HttpResponse:
    doc = Document()
    doc.add_heading(f"Notes on {topic}", 0)

    for line in (notes_text or "").split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("NOTE:"):
            para = doc.add_paragraph(line)
            # ensure a run exists
            run = para.runs[0] if para.runs else para.add_run(line)
            run.bold = True
            run.font.color.rgb = DocxRGBColor(255, 0, 0)
        else:
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    response = HttpResponse(
        buffer,
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    response['Content-Disposition'] = f'attachment; filename="{topic}.docx"'
    return response


# ---------------- PPT ----------------
def generate_ppt(topic, notes_input) -> HttpResponse:
    """
    notes_input: either a structured list of dicts OR a plain text string.
    Structured list item example:
      {'type': 'heading'/'subheading'/'paragraph'/'note'/'list', 'text': '...', 'level': 0}
    """
    prs = Presentation()
    slide = None
    tf = None

    def create_new_slide(title_text):
        layout = prs.slide_layouts[1]  # Title + Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title_text
        tx = s.placeholders[1].text_frame
        tx.clear()
        return s, tx

    # CASE A: structured list supplied (preferred)
    if isinstance(notes_input, list):
        notes = notes_input
        for item in notes:
            typ = item.get('type')
            text = (item.get('text') or "").strip()
            if not text:
                continue

            if typ == 'heading':
                slide, tf = create_new_slide(text)
                continue

            # ensure a slide exists
            if slide is None:
                slide, tf = create_new_slide(f"Notes on {topic}")

            if typ in ('subheading', 'subsubheading'):
                p = tf.add_paragraph()
                p.text = text
                p.font.bold = True
                p.font.size = Pt(16 if typ == 'subheading' else 14)
            elif typ == 'paragraph':
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(12)
            elif typ == 'note':
                p = tf.add_paragraph()
                p.text = f"NOTE: {text}"
                p.font.bold = True
                try:
                    p.font.color.rgb = PptxRGBColor(255, 0, 0)
                except Exception:
                    # ignore color setting if any unexpected error
                    pass
                p.font.size = Pt(12)
            elif typ == 'list':
                p = tf.add_paragraph()
                p.text = text
                # support optional nesting level
                level = item.get('level', 0) or 0
                try:
                    p.level = int(level)
                except Exception:
                    p.level = 0
                p.font.size = Pt(12)
        # if we never created a slide (empty structured input) create default
        if slide is None:
            slide, tf = create_new_slide(f"Notes on {topic}")

    else:
        # CASE B: fallback - plain text string
        notes_text = (notes_input or "")
        lines = [ln.strip() for ln in notes_text.splitlines() if ln.strip()]
        for line in lines:
            # heuristic to treat a line as heading:
            is_heading = False
            # - explicit prefix '## ' as heading
            if line.startswith("## "):
                is_heading = True
                heading_text = line[3:].strip()
            # - all uppercase and reasonably short => heading (e.g., "INTRODUCTION")
            elif line.isupper() and len(line) > 2 and len(line) < 120:
                is_heading = True
                heading_text = line
            # - short line ending with ':' likely a heading
            elif line.endswith(':') and len(line) < 120:
                is_heading = True
                heading_text = line.rstrip(':').strip()
            else:
                heading_text = None

            if is_heading:
                slide, tf = create_new_slide(heading_text)
                continue

            # ensure a slide exists
            if slide is None:
                slide, tf = create_new_slide(f"Notes on {topic}")

            # now add content to current slide
            if line.startswith("NOTE:"):
                p = tf.add_paragraph()
                p.text = line
                p.font.bold = True
                try:
                    p.font.color.rgb = PptxRGBColor(255, 0, 0)
                except Exception:
                    pass
                p.font.size = Pt(12)
            elif line.startswith("•") or line.startswith("-"):
                p = tf.add_paragraph()
                p.text = line
                p.level = 1
                p.font.size = Pt(12)
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)

    # Save and return
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    response = HttpResponse(
        buffer,
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename="{topic}.pptx"'
    return response
